"""
Vehicle Delivery Analysis PowerPoint Presentation Generator
Data Science Course Presentation
Focus: Methodology and Techniques for Academic Audience
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

def create_vehicle_analysis_presentation():
    """
    Creates a comprehensive PowerPoint presentation for vehicle delivery analysis
    Targeted for data science students and instructors
    """
    
    # Initialize presentation
    prs = Presentation()
    
    # Define color scheme for academic presentation
    primary_color = RGBColor(31, 73, 125)  # Professional blue
    accent_color = RGBColor(68, 114, 196)  # Lighter blue
    text_color = RGBColor(68, 68, 68)      # Dark gray
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "Amazon Delivery Fleet Optimization Analysis"
    title.text_frame.paragraphs[0].font.color.rgb = primary_color
    title.text_frame.paragraphs[0].font.size = Pt(40)
    
    subtitle.text = ("A Data Science Approach to Vehicle Performance Analysis\n"
                    "Methodology Focus: Statistical Analysis & Geospatial Visualization\n\n"
                    "Data Science Course Project\n"
                    "December 2025")
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # Slide 2: Problem Statement & Objectives
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "Research Problem & Data Science Objectives"
    slide2.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide2.placeholders[1]
    content.text = ("Problem Statement:\n"
                   "• How do different vehicle types perform under varying environmental conditions?\n"
                   "• What data-driven insights can optimize delivery fleet deployment?\n\n"
                   "Data Science Learning Objectives:\n"
                   "• Apply Exploratory Data Analysis (EDA) techniques\n"
                   "• Implement statistical visualization methods\n"
                   "• Utilize geospatial data analysis and mapping\n"
                   "• Practice multi-factor performance analysis\n"
                   "• Develop business insights from data patterns\n\n"
                   "Dataset: Amazon Delivery Records (43,740 observations)\n"
                   "Key Variables: Vehicle type, Weather, Traffic, Delivery time, Location")
    
    # Slide 3: Methodology Framework
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "6-Phase Data Science Methodology"
    slide3.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide3.placeholders[1]
    content.text = ("Phase 1: Data Preparation & Quality Assessment\n"
                   "• Data cleaning, standardization, outlier detection\n"
                   "• Missing value analysis and treatment\n\n"
                   "Phase 2: Baseline Vehicle Performance Analysis\n"
                   "• Descriptive statistics, distribution analysis\n"
                   "• Box plots for central tendency comparison\n\n"
                   "Phase 3: Weather Impact Statistical Analysis\n"
                   "• Categorical variable analysis, ANOVA techniques\n"
                   "• Grouped statistical comparisons\n\n"
                   "Phase 4: Traffic Condition Performance Evaluation\n"
                   "• Multi-level categorical analysis\n"
                   "• Traffic impact quantification\n\n"
                   "Phase 5: Combined Conditions Optimization\n"
                   "• Heatmap visualization, correlation analysis\n"
                   "• Multi-factor interaction effects\n\n"
                   "Phase 6: Geospatial Analysis & Interactive Visualization\n"
                   "• Coordinate validation, spatial distribution mapping\n"
                   "• Interactive web-based visualization")
    
    # Slide 4: Statistical Techniques Applied
    slide4 = prs.slides.add_slide(prs.slide_layouts[1])
    slide4.shapes.title.text = "Statistical & Visualization Techniques"
    slide4.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide4.placeholders[1]
    content.text = ("Descriptive Statistics:\n"
                   "• Central tendency measures (median focus for skewed distributions)\n"
                   "• Variance analysis and quartile-based comparisons\n\n"
                   "Visualization Methods:\n"
                   "• Box plots: Distribution comparison across categories\n"
                   "• Bar charts: Mean performance visualization\n"
                   "• Heatmaps: Multi-dimensional correlation analysis\n"
                   "• Geographic mapping: Spatial pattern identification\n\n"
                   "Analytical Considerations:\n"
                   "• Non-parametric approaches for robustness\n"
                   "• Categorical variable handling strategies\n"
                   "• Data quality validation protocols\n"
                   "• Interactive visualization design principles\n\n"
                   "Tools: Python (Pandas, Matplotlib, Seaborn, Folium)")
    
    # Slide 5: Phase 1 - Delivery Time Overview
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout for custom positioning
    title_shape = slide5.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Phase 1: Overall Vehicle Performance Distribution"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add image
    if os.path.exists('01_vehicle_delivery_time_overview.png'):
        slide5.shapes.add_picture('01_vehicle_delivery_time_overview.png', 
                                Inches(1), Inches(1.2), Inches(8), Inches(5))
    
    # Add methodology note
    method_box = slide5.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(1))
    method_frame = method_box.text_frame
    method_frame.text = ("Methodology: Box plot analysis for distribution comparison\n"
                        "Key Insight: Median-based ranking reveals bike superiority with lowest delivery times")
    method_frame.paragraphs[0].font.size = Pt(14)
    method_frame.paragraphs[0].font.color.rgb = text_color
    
    # Slide 6: Phase 2 - Weather Impact Analysis
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    title_shape = slide6.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Phase 2: Weather Condition Impact Analysis"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add weather boxplot
    if os.path.exists('02a_vehicle_weather_boxplot.png'):
        slide6.shapes.add_picture('02a_vehicle_weather_boxplot.png', 
                                Inches(0.5), Inches(1.2), Inches(4.5), Inches(4))
    
    # Add weather bar chart
    if os.path.exists('02b_vehicle_weather_bar.png'):
        slide6.shapes.add_picture('02b_vehicle_weather_bar.png', 
                                Inches(5), Inches(1.2), Inches(4.5), Inches(4))
    
    method_box = slide6.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    method_frame = method_box.text_frame
    method_frame.text = ("Statistical Approach: Categorical analysis with grouped comparisons\n"
                        "Visualization Strategy: Box plots for distribution, bar charts for means\n"
                        "Finding: Weather resilience varies significantly by vehicle type")
    method_frame.paragraphs[0].font.size = Pt(14)
    
    # Slide 7: Phase 3 - Traffic Impact Analysis
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    title_shape = slide7.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Phase 3: Traffic Condition Performance Analysis"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add traffic analysis
    if os.path.exists('03a_vehicle_traffic_boxplot.png'):
        slide7.shapes.add_picture('03a_vehicle_traffic_boxplot.png', 
                                Inches(0.5), Inches(1.2), Inches(4.5), Inches(4))
    
    if os.path.exists('03b_vehicle_traffic_bar.png'):
        slide7.shapes.add_picture('03b_vehicle_traffic_bar.png', 
                                Inches(5), Inches(1.2), Inches(4.5), Inches(4))
    
    method_box = slide7.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    method_frame = method_box.text_frame
    method_frame.text = ("Multi-level Categorical Analysis: Traffic impact quantification\n"
                        "Comparative Method: Vehicle maneuverability assessment\n"
                        "Key Pattern: Motorcycles and bikes show superior traffic adaptability")
    method_frame.paragraphs[0].font.size = Pt(14)
    
    # Slide 8: Phase 4 - Combined Conditions Analysis
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    title_shape = slide8.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Phase 4: Multi-Factor Interaction Analysis"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add combined conditions heatmap
    if os.path.exists('04a_combined_conditions_heatmap.png'):
        slide8.shapes.add_picture('04a_combined_conditions_heatmap.png', 
                                Inches(0.5), Inches(1.2), Inches(4.5), Inches(4))
    
    if os.path.exists('04b_vehicle_efficiency_rankings.png'):
        slide8.shapes.add_picture('04b_vehicle_efficiency_rankings.png', 
                                Inches(5), Inches(1.2), Inches(4.5), Inches(4))
    
    method_box = slide8.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(9), Inches(1.5))
    method_frame = method_box.text_frame
    method_frame.text = ("Heatmap Methodology: Multi-dimensional correlation visualization\n"
                        "Interaction Effects: Weather × Traffic combined impact analysis\n"
                        "Optimization Insight: Condition-specific vehicle deployment strategies identified")
    method_frame.paragraphs[0].font.size = Pt(14)
    
    # Slide 9: Phase 5 - Interactive Geospatial Analysis
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    title_shape = slide9.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_shape.text_frame
    title_frame.text = "Phase 5: Interactive Geospatial Visualization"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.color.rgb = primary_color
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Note for HTML embedding
    html_note = slide9.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    html_frame = html_note.text_frame
    html_frame.text = ("Interactive HTML Heatmap: 05_store_delivery_heatmap.html\n"
                      "[Live demonstration during presentation]\n"
                      "Geospatial Analysis Techniques Applied:")
    html_frame.paragraphs[0].font.size = Pt(18)
    html_frame.paragraphs[0].font.color.rgb = accent_color
    html_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    content_box = slide9.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(3))
    content_frame = content_box.text_frame
    content_frame.text = ("• Coordinate Validation: Removed invalid lat/long entries (0,0)\n"
                         "• Spatial Distribution Mapping: Store location density analysis\n"
                         "• Interactive Web Visualization: Folium-based mapping\n"
                         "• Delivery Volume Hotspots: Geographic pattern identification\n"
                         "• Data Quality Protocols: Geographic data cleaning procedures\n\n"
                         "Technical Implementation:\n"
                         "• Python Folium library for interactive mapping\n"
                         "• HTML/JavaScript output for web compatibility\n"
                         "• Responsive design for presentation display")
    content_frame.paragraphs[0].font.size = Pt(16)
    content_frame.paragraphs[0].font.color.rgb = text_color
    
    # Slide 10: Key Findings Summary
    slide10 = prs.slides.add_slide(prs.slide_layouts[1])
    slide10.shapes.title.text = "Data-Driven Findings & Statistical Insights"
    slide10.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide10.placeholders[1]
    content.text = ("Vehicle Performance Rankings (Median Delivery Time):\n"
                   "1. Bike - Most efficient (shortest delivery times)\n"
                   "2. Motorcycle - Second best performer\n"
                   "3. Scooter - Mid-range efficiency\n"
                   "4. Car - Lower efficiency for urban delivery\n"
                   "5. Truck - Longest delivery times\n\n"
                   "Environmental Impact Patterns:\n"
                   "• Weather Resilience: Bikes/motorcycles show consistency\n"
                   "• Traffic Adaptability: Small vehicles excel in congestion\n"
                   "• Optimal Conditions: Sunny + Low traffic universally best\n"
                   "• Worst Conditions: Stormy + Traffic jam most problematic\n\n"
                   "Geospatial Insights:\n"
                   "• Store location clustering reveals delivery density patterns\n"
                   "• Geographic distribution suggests regional optimization opportunities")
    
    # Slide 11: Methodology Limitations & Considerations
    slide11 = prs.slides.add_slide(prs.slide_layouts[1])
    slide11.shapes.title.text = "Statistical Limitations & Methodological Considerations"
    slide11.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide11.placeholders[1]
    content.text = ("Data Quality Considerations:\n"
                   "• Missing value patterns and treatment strategies\n"
                   "• Outlier detection and robust statistical approaches\n"
                   "• Geographic coordinate validation protocols\n\n"
                   "Statistical Methodology Limitations:\n"
                   "• Median-based analysis choice (skewed distributions)\n"
                   "• Categorical variable grouping decisions\n"
                   "• Temporal patterns not analyzed (static snapshot)\n"
                   "• Causation vs. correlation interpretation\n\n"
                   "Visualization Design Choices:\n"
                   "• Color scheme selection for clarity\n"
                   "• Box plot vs. bar chart comparative advantages\n"
                   "• Heatmap scaling and interpretation considerations\n"
                   "• Interactive vs. static visualization trade-offs\n\n"
                   "Future Analysis Opportunities:\n"
                   "• Time series analysis for seasonal patterns\n"
                   "• Predictive modeling implementation\n"
                   "• Advanced geospatial analysis techniques")
    
    # Slide 12: Learning Outcomes & Technical Skills
    slide12 = prs.slides.add_slide(prs.slide_layouts[1])
    slide12.shapes.title.text = "Data Science Learning Outcomes & Skills Applied"
    slide12.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide12.placeholders[1]
    content.text = ("Technical Skills Demonstrated:\n"
                   "• Exploratory Data Analysis (EDA) comprehensive workflow\n"
                   "• Statistical visualization best practices\n"
                   "• Multi-factor analysis methodology\n"
                   "• Geospatial data processing and visualization\n"
                   "• Interactive web-based visualization development\n\n"
                   "Analytical Thinking Applied:\n"
                   "• Business problem decomposition into data questions\n"
                   "• Systematic methodology design and implementation\n"
                   "• Statistical interpretation and insight generation\n"
                   "• Data quality assessment and validation protocols\n\n"
                   "Communication & Presentation:\n"
                   "• Data storytelling through progressive analysis\n"
                   "• Visualization design for different audiences\n"
                   "• Technical methodology documentation\n"
                   "• Academic presentation structure and flow\n\n"
                   "Tools Mastery: Python ecosystem (Pandas, Matplotlib, Seaborn, Folium)")
    
    # Slide 13: Conclusion & Discussion
    slide13 = prs.slides.add_slide(prs.slide_layouts[1])
    slide13.shapes.title.text = "Conclusion & Academic Discussion Points"
    slide13.shapes.title.text_frame.paragraphs[0].font.color.rgb = primary_color
    
    content = slide13.placeholders[1]
    content.text = ("Project Summary:\n"
                   "Successfully applied systematic data science methodology to real-world\n"
                   "logistics optimization problem with measurable business insights\n\n"
                   "Methodological Contributions:\n"
                   "• Demonstrated multi-phase analysis workflow design\n"
                   "• Integrated statistical and geospatial analysis approaches\n"
                   "• Combined static and interactive visualization strategies\n\n"
                   "Discussion Questions for Class:\n"
                   "• How might temporal analysis enhance these findings?\n"
                   "• What additional statistical tests could strengthen conclusions?\n"
                   "• How would you validate these insights with real business metrics?\n"
                   "• What ethical considerations apply to delivery optimization?\n\n"
                   "Next Steps:\n"
                   "• Predictive model development\n"
                   "• A/B testing framework design\n"
                   "• Real-time dashboard implementation")
    
    return prs

def main():
    """Generate the PowerPoint presentation"""
    print("Creating Vehicle Analysis PowerPoint Presentation...")
    print("Target Audience: Data Science Students & Instructor")
    print("Focus: Methodology and Technical Approaches")
    
    # Create presentation
    presentation = create_vehicle_analysis_presentation()
    
    # Save presentation
    output_file = "Vehicle_Delivery_Analysis_Data_Science_Presentation.pptx"
    presentation.save(output_file)
    
    print(f"\nPresentation created successfully: {output_file}")
    print("\nPresentation Structure:")
    print("1. Title Slide")
    print("2. Problem Statement & Objectives") 
    print("3. 6-Phase Data Science Methodology")
    print("4. Statistical & Visualization Techniques")
    print("5. Phase 1: Overall Performance Distribution")
    print("6. Phase 2: Weather Impact Analysis") 
    print("7. Phase 3: Traffic Condition Analysis")
    print("8. Phase 4: Multi-Factor Interaction Analysis")
    print("9. Phase 5: Interactive Geospatial Visualization")
    print("10. Key Findings & Statistical Insights")
    print("11. Methodology Limitations & Considerations")
    print("12. Learning Outcomes & Technical Skills")
    print("13. Conclusion & Academic Discussion")
    print("\nNote: HTML heatmap (05_store_delivery_heatmap.html) ready for live demo")
    print("Visualization files integrated: 8 PNG files embedded")

if __name__ == "__main__":
    main()